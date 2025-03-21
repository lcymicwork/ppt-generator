from flask import Flask, request, jsonify, send_file
from pptx import Presentation
from pptx.util import Inches
import io
import os
import requests
from openai import OpenAI
import json  # Import the json module

app = Flask(__name__)

OPENROUTER_API_KEY = os.environ.get("OPENROUTER_API_KEY")
SILICONFLOW_API_URL = "https://api.siliconflow.cn/v1/images/generations"
OPENROUTER_API_URL = "https://openrouter.ai/api/v1"
SILICONFLOW_API_TOKEN = "sk-mdkwkbklsycmjxjhstkxpfgsuzmfnwxwsqrnbmqlhxljitsq"

@app.route('/test-siliconflow', methods=['POST'])
def test_siliconflow():
    url = "https://api.siliconflow.cn/v1/images/generations"
    headers = {
        "Authorization": f"Bearer {SILICONFLOW_API_TOKEN}",
        "Content-Type": "application/json"
    }
    try:
        response = requests.get(url, headers=headers)
        response.raise_for_status()  # Raise HTTPError for bad responses (4xx or 5xx)
        return jsonify({'success': True})
    except requests.exceptions.RequestException as e:
        return jsonify({'success': False, 'error': str(e)})

@app.route('/test-openrouter', methods=['POST'])
def test_openrouter():
    data = request.get_json()
    api_key = data.get('apiKey')

    if not api_key:
        return jsonify({'success': False, 'error': 'API key is required'}), 400

    try:
        url = "https://openrouter.ai/api/v1/chat/completions"
        headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
            "HTTP-Referer": "http://localhost:5000",  # Optional
            "X-Title": "AI PPT Generator",  # Optional
        }
        payload = json.dumps({
            "model": "google/gemini-2.0-flash-001",
            "messages": [{"role": "user", "content": [{"type": "text", "text": "Say hi"}]}],
        })

        response = requests.post(url, headers=headers, data=payload)
        response.raise_for_status()  # Raise HTTPError for bad responses

        # Attempt to parse the JSON response
        try:
            response_json = response.json()
            return jsonify({'success': True, 'response': response_json})
        except json.JSONDecodeError as e:
            print(f"JSONDecodeError: {e}, Response Text: {response.text}")
            return jsonify({'success': False, 'error': f"Invalid JSON response: {str(e)}. Raw response: {response.text}"}), 500

    except requests.exceptions.RequestException as e:
        print(f"OpenRouter API error: {e}")
        return jsonify({'success': False, 'error': str(e)}), 500
    except Exception as e:
        print(f"General error: {e}")
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/generate-content-only', methods=['POST'])
def generate_content_only():
    data = request.get_json()
    prompt = data['prompt']
    api_key = data['apiKey']
    openrouter_llm = data['openrouterLlm']

    content = generate_content(prompt, api_key, openrouter_llm)
    return jsonify(content)

@app.route('/generate-ppt', methods=['POST'])
def generate_ppt():
    data = request.get_json()
    prompt = data['prompt']
    api_key = data['apiKey']
    openrouter_llm = data['openrouterLlm']
    siliconflow_llm = data['siliconflowLlm']

    # If the prompt is empty, generate content
    if not prompt:
        content = generate_content(prompt, api_key, openrouter_llm)
    else:
        # Parse the content from the textarea
        content = []
        slides = prompt.split('\n\n')
        for slide in slides:
            if slide.strip():
                lines = slide.split('\n')
                title = lines[0].split(': ')[1] if len(lines) > 0 and 'Title:' in lines[0] else 'No Title'
                content_text = lines[1].split(': ')[1] if len(lines) > 1 and 'Content:' in lines[1] else 'No Content'
                content.append({'title': title, 'content': content_text})

    # Create a PowerPoint presentation
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "AI Generated PPT"
    subtitle.text = prompt

    # Add slides based on the generated content
    for slide_content in content:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        body_shape = slide.shapes.placeholders[1]

        tf = body_shape.text_frame
        title.text = slide_content['title']
        tf.text = slide_content['content']

        # Generate image using siliconflow.cn
        image_url = generate_image(slide_content['title'], siliconflow_llm)
        if image_url:
            try:
                image_response = requests.get(image_url, stream=True)
                image_response.raise_for_status()
                image_data = image_response.content
                image_stream = io.BytesIO(image_data)

                left = top = Inches(2.0)
                pic = slide.shapes.add_picture(image_stream, left, top, width=Inches(6.0))
            except requests.exceptions.RequestException as e:
                print(f"Error downloading image: {e}")

    # Save the presentation to a byte stream
    pptx_stream = io.BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)

    # Return the presentation as a downloadable file
    return send_file(
        pptx_stream,
        as_attachment=True,
        download_name='presentation.pptx',
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )

def generate_content(prompt, api_key, openrouter_llm):
    client = OpenAI(
        base_url=OPENROUTER_API_URL,
        api_key=api_key,
    )

    completion = client.chat.completions.create(
        extra_headers={
            "HTTP-Referer": "http://localhost:5000",  # Optional. Site URL for rankings on openrouter.ai.
            "X-Title": "AI PPT Generator",  # Optional. Site title for rankings on openrouter.ai.
        },
        extra_body={},
        model=openrouter_llm or "google/gemini-2.0-flash-001",
        messages=[
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": prompt
                    }
                ]
            }
        ]
    )
    return [{'title': 'Slide 1', 'content': completion.choices[0].message.content}]

def generate_image(prompt, siliconflow_llm):
    url = "https://api.siliconflow.cn/v1/images/generations"

    payload = {
        "model": siliconflow_llm or "Kwai-Kolors/Kolors",
        "prompt": prompt,
        "negative_prompt": "",
        "image_size": "1024x1024",
        "batch_size": 1,
        "seed": 4999999999,
        "num_inference_steps": 20,
        "guidance_scale": 7.5,
        "image": ""
    }
    headers = {
        "Authorization": f"Bearer {SILICONFLOW_API_TOKEN}",
        "Content-Type": "application/json"
    }

    response = requests.request("POST", url, json=payload, headers=headers)

    try:
        return response.json()['images'][0]
    except:
        return None

if __name__ == '__main__':
    app.run(debug=True)
